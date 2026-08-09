import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGBColor, Inches as DocxInches

# 💡 PPTX / DOCX 자동 내보내기 도우미 모듈 신설
# LLM이 출력한 마크다운/슬라이드 구조의 텍스트를 파싱하여 진짜 파워포인트 및 워드 파일로 렌더링합니다.

def create_styled_pptx(llm_output: str, output_path: str = "Tech_GPT_Presentation.pptx") -> str:
    """
    LLM의 슬라이드 마크다운 텍스트를 깔끔한 디자인 테마가 적용된 실제 .pptx 슬라이드로 변환
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 와이드스크린 사양
    prs.slide_height = Inches(7.5)

    # 마크다운 구분선(---) 기준 슬라이드 분할
    raw_slides = re.split(r'---\s*', llm_output)

    for slide_raw in raw_slides:
        if not slide_raw.strip():
            continue

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # 1. 디자인 상단 헤더 배너 배경 (다크 네이비 테마)
        header_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.333), Inches(1.2)  # 1 = MSO_SHAPE.RECTANGLE
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(24, 43, 73)  # Tech-Navy
        header_shape.line.fill.background()

        # 2. 제목 추출 및 파싱
        title_match = re.search(r'###\s*🖼️\s*\[?(Slide\s*\d+.*?)\]?', slide_raw)
        if not title_match:
            title_match = re.search(r'###\s*(.*)', slide_raw)
        
        slide_title = title_match.group(1).strip() if title_match else "Tech-GPT 기술 분석 슬라이드"

        # 헤더 제목 텍스트 박스
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_title
        p_title.font.bold = True
        p_title.font.size = Pt(22)
        p_title.font.color.rgb = RGBColor(255, 255, 255)

        # 3. 본문 콘텐츠 텍스트 박스
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True

        lines = slide_raw.strip().split("\n")
        first_line = True

        for line in lines:
            line_str = line.strip()
            if not line_str or line_str.startswith("###"):
                continue

            p = tf_body.add_paragraph() if not first_line else tf_body.paragraphs[0]
            first_line = False

            # 불릿 및 스크립트 강조 정제
            clean_text = line_str.replace("**", "").replace("*", "").replace("#", "").strip()

            if line_str.startswith("- **핵심 요약**") or line_str.startswith("- **핵심"):
                p.text = f"💡 {clean_text}"
                p.font.bold = True
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(0, 102, 204)
            elif line_str.startswith("> 💡") or "발표자 스크립트" in line_str:
                p.text = f"🗣️ {clean_text}"
                p.font.italic = True
                p.font.size = Pt(13)
                p.font.color.rgb = RGBColor(100, 100, 100)
            else:
                p.text = f"• {clean_text}"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(40, 40, 40)

    prs.save(output_path)
    return output_path


def create_styled_docx(llm_output: str, output_path: str = "Tech_GPT_Report.docx") -> str:
    """
    LLM의 보고서 마크다운 텍스트를 양식이 갖춰진 실제 .docx MS 워드 문서로 변환
    """
    doc = Document()

    # 여백 설정
    for section in doc.sections:
        section.top_margin = DocxInches(1)
        section.bottom_margin = DocxInches(1)
        section.left_margin = DocxInches(1)
        section.right_margin = DocxInches(1)

    lines = llm_output.split("\n")

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        if stripped.startswith("# "):
            h = doc.add_heading(level=1)
            run = h.add_run(stripped.replace("# ", ""))
            run.font.size = DocxPt(20)
            run.font.bold = True
            run.font.color.rgb = DocxRGBColor(24, 43, 73)
        elif stripped.startswith("## "):
            h = doc.add_heading(level=2)
            run = h.add_run(stripped.replace("## ", ""))
            run.font.size = DocxPt(15)
            run.font.bold = True
            run.font.color.rgb = DocxRGBColor(0, 102, 204)
        elif stripped.startswith("### "):
            h = doc.add_heading(level=3)
            run = h.add_run(stripped.replace("### ", ""))
            run.font.size = DocxPt(13)
            run.font.bold = True
        elif stripped.startswith("- ") or stripped.startswith("* "):
            p = doc.add_paragraph(style='List Bullet')
            clean_text = stripped.lstrip("-* ").replace("**", "")
            p.add_run(clean_text)
        else:
            p = doc.add_paragraph()
            clean_text = stripped.replace("**", "")
            p.add_run(clean_text)

    doc.save(output_path)
    return output_path
# =====수정사항 종료=====